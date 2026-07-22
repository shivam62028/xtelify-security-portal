"""
Generate PowerPoint Presentation for Xtelify Security Portal
Run: pip install python-pptx && python create_presentation.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Create presentation with 16:9 aspect ratio
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
DARK_BLUE = RGBColor(15, 23, 42)
LIGHT_BLUE = RGBColor(14, 165, 233)
PURPLE = RGBColor(139, 92, 246)
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(226, 232, 240)
GREEN = RGBColor(16, 185, 129)


def add_title_slide(title, subtitle):
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Background
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BLUE
    background.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(28)
    p.font.color.rgb = LIGHT_BLUE
    p.alignment = PP_ALIGN.CENTER

    return slide


def add_content_slide(title, bullets, code_block=None):
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Background
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BLUE
    background.line.fill.background()

    # Title bar
    title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = RGBColor(30, 41, 59)
    title_bar.line.fill.background()

    # Title text
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = LIGHT_BLUE

    # Bullets
    if bullets:
        y_pos = Inches(1.5)
        if code_block:
            y_pos = Inches(1.4)

        bullet_box = slide.shapes.add_textbox(Inches(0.7), y_pos, Inches(12), Inches(3.5) if code_block else Inches(5.5))
        tf = bullet_box.text_frame
        tf.word_wrap = True

        for i, bullet in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = "  " + bullet
            p.font.size = Pt(20) if code_block else Pt(22)
            p.font.color.rgb = LIGHT_GRAY
            p.space_after = Pt(8)

    # Code block if provided
    if code_block:
        code_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(4.3), Inches(12.333), Inches(2.7))
        code_box.fill.solid()
        code_box.fill.fore_color.rgb = RGBColor(30, 41, 59)
        code_box.line.color.rgb = RGBColor(51, 65, 85)

        code_text = slide.shapes.add_textbox(Inches(0.7), Inches(4.5), Inches(12), Inches(2.4))
        tf = code_text.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = code_block
        p.font.size = Pt(14)
        p.font.name = "Consolas"
        p.font.color.rgb = RGBColor(165, 243, 252)

    return slide


def add_two_column_slide(title, left_title, left_items, right_title, right_items):
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Background
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BLUE
    background.line.fill.background()

    # Title bar
    title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = RGBColor(30, 41, 59)
    title_bar.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = LIGHT_BLUE

    # Left column box
    left_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(6), Inches(5.5))
    left_box.fill.solid()
    left_box.fill.fore_color.rgb = RGBColor(30, 41, 59)
    left_box.line.color.rgb = RGBColor(51, 65, 85)

    # Left title
    lt_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.7), Inches(5.5), Inches(0.5))
    tf = lt_box.text_frame
    p = tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = PURPLE

    # Left items
    li_box = slide.shapes.add_textbox(Inches(0.7), Inches(2.3), Inches(5.5), Inches(4.5))
    tf = li_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(left_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "  " + item
        p.font.size = Pt(17)
        p.font.color.rgb = LIGHT_GRAY
        p.space_after = Pt(6)

    # Right column box
    right_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.833), Inches(1.5), Inches(6), Inches(5.5))
    right_box.fill.solid()
    right_box.fill.fore_color.rgb = RGBColor(30, 41, 59)
    right_box.line.color.rgb = RGBColor(51, 65, 85)

    # Right title
    rt_box = slide.shapes.add_textbox(Inches(7), Inches(1.7), Inches(5.5), Inches(0.5))
    tf = rt_box.text_frame
    p = tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = PURPLE

    # Right items
    ri_box = slide.shapes.add_textbox(Inches(7), Inches(2.3), Inches(5.5), Inches(4.5))
    tf = ri_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(right_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "  " + item
        p.font.size = Pt(17)
        p.font.color.rgb = LIGHT_GRAY
        p.space_after = Pt(6)

    return slide


def add_architecture_slide():
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Background
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BLUE
    background.line.fill.background()

    # Title bar
    title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = RGBColor(30, 41, 59)
    title_bar.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "System Architecture"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = LIGHT_BLUE

    # Browser box
    browser = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(1.5), Inches(11.333), Inches(2))
    browser.fill.solid()
    browser.fill.fore_color.rgb = RGBColor(30, 41, 59)
    browser.line.color.rgb = LIGHT_BLUE

    browser_text = slide.shapes.add_textbox(Inches(1.2), Inches(1.7), Inches(11), Inches(1.8))
    tf = browser_text.text_frame
    p = tf.paragraphs[0]
    p.text = "BROWSER (React + TypeScript)"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = LIGHT_BLUE
    p = tf.add_paragraph()
    p.text = "Dashboard  |  Charts (Recharts)  |  Upload Modal  |  Export (XLSX/PDF)"
    p.font.size = Pt(16)
    p.font.color.rgb = LIGHT_GRAY

    # Arrow
    arrow = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(6.2), Inches(3.6), Inches(0.8), Inches(0.7))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = LIGHT_BLUE
    arrow.line.fill.background()

    # Arrow label
    arrow_label = slide.shapes.add_textbox(Inches(7.2), Inches(3.7), Inches(2), Inches(0.5))
    tf = arrow_label.text_frame
    p = tf.paragraphs[0]
    p.text = "REST API (HTTP)"
    p.font.size = Pt(14)
    p.font.color.rgb = LIGHT_GRAY

    # Backend box
    backend = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(4.4), Inches(11.333), Inches(2.6))
    backend.fill.solid()
    backend.fill.fore_color.rgb = RGBColor(30, 41, 59)
    backend.line.color.rgb = PURPLE

    backend_text = slide.shapes.add_textbox(Inches(1.2), Inches(4.6), Inches(11), Inches(2.3))
    tf = backend_text.text_frame
    p = tf.paragraphs[0]
    p.text = "BACKEND (FastAPI + Python)"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = PURPLE
    p = tf.add_paragraph()
    p.text = "/api/upload-report  |  /api/db  |  /api/leaderboard"
    p.font.size = Pt(16)
    p.font.color.rgb = LIGHT_GRAY
    p = tf.add_paragraph()
    p.text = "Smart Worksheet Detection (openpyxl)  |  Data Processing (Pandas)"
    p.font.size = Pt(16)
    p.font.color.rgb = LIGHT_GRAY
    p = tf.add_paragraph()
    p.text = "Storage: xtelify_db.json"
    p.font.size = Pt(16)
    p.font.color.rgb = GREEN

    return slide


# ============ CREATE SLIDES ============

# Slide 1: Title
add_title_slide(
    "Xtelify Security Portal",
    "Enterprise Vulnerability Management Dashboard"
)

# Slide 2: Overview
add_content_slide(
    "Project Overview",
    [
        "Web-based vulnerability management platform for enterprise security teams",
        "Upload and process vulnerability scan reports (Excel, CSV, JSON)",
        "Smart worksheet detection for multi-sheet Excel files",
        "Real-time vulnerability tracking and status management",
        "Team performance leaderboard with MTTR metrics",
        "Export capabilities (Excel, PDF)",
        "Role-based access control (Admin / Viewer)"
    ]
)

# Slide 3: Tech Stack
add_two_column_slide(
    "Tech Stack",
    "Frontend",
    [
        "React 19 - UI Framework",
        "TypeScript - Type Safety",
        "Vite 8 - Build Tool",
        "TailwindCSS 4 - Styling",
        "Recharts - Data Visualization",
        "Lucide React - Icons",
        "XLSX (SheetJS) - Excel Export",
        "jsPDF - PDF Generation"
    ],
    "Backend",
    [
        "Python 3.x - Language",
        "FastAPI - REST API Framework",
        "Pandas - Data Processing",
        "openpyxl - Excel Inspection",
        "httpx - Async HTTP Client",
        "LangSmith - AI Observability",
        "JSON File - Database Storage",
        "Docker + Nginx - Deployment"
    ]
)

# Slide 4: Architecture
add_architecture_slide()

# Slide 5: API Endpoints
add_content_slide(
    "API Endpoints",
    [
        "GET  /api/db                    - Fetch all vulnerability records",
        "POST /api/db                    - Add new vulnerability items",
        "DELETE /api/db                  - Delete dataset by UploadBatch",
        "POST /api/upload-report         - Upload & process Excel/CSV/JSON",
        "POST /api/upload-report-with-sheet - Upload with manual sheet selection",
        "GET  /api/leaderboard           - Get team performance rankings",
        "GET  /{path}                    - Serve frontend static files"
    ]
)

# Slide 6: Upload Flow
add_content_slide(
    "Data Flow: Upload Process",
    [
        "1. User selects Excel file in browser",
        "2. Frontend sends FormData to /api/upload-report",
        "3. Backend scans all worksheets using openpyxl (read-only mode)",
        "4. Smart detection scores each sheet and selects the best one",
        "5. Header row detected by scanning first 15 rows",
        "6. Pandas reads selected sheet from header row",
        "7. Column mapping applied (ID -> IssueID, CVSSSeverity -> Severity)",
        "8. Records normalized and saved to xtelify_db.json",
        "9. Frontend reloads to display new data"
    ]
)

# Slide 7: Smart Worksheet Detection
add_two_column_slide(
    "Smart Worksheet Detection",
    "Positive Scoring",
    [
        "+3 pts: ID, Name, Severity, FindingStatus",
        "+2 pts: Score, CVSSSeverity, WizURL",
        "+1 pt: Other expected columns",
        "+5 pts: More than 100 data rows",
        "+3 pts: More than 50 data rows"
    ],
    "Negative Scoring",
    [
        "-5 pts: 'Grand Total' in sheet name",
        "-5 pts: 'Count of' in sheet name",
        "-5 pts: 'Pivot' in sheet name",
        "-3 pts: Fewer than 20 rows",
        "",
        "Auto-selects highest scoring sheet!"
    ]
)

# Slide 8: Column Mapping
add_content_slide(
    "Column Auto-Mapping",
    [
        "Supports multiple column naming conventions from different scanners:"
    ],
    "ID / IssueID / VulnID / CVE           ->  IssueID\nCVSSSeverity / VendorSeverity / Risk  ->  Severity\nFirstDetected / DetectedDate           ->  DiscoveredDate\nWizURL / Link / Reference              ->  ReferenceLinks\nRemediation / Resolution / Fix         ->  RecommendedAction"
)

# Slide 9: Vulnerability Grouping
add_content_slide(
    "Vulnerability Grouping",
    [
        "Frontend groups vulnerabilities by CVE/ID",
        "Multiple assets affected by same vulnerability shown together",
        "Expandable rows to see affected assets",
        "Highest severity from group displayed",
        "Status aggregated (Open if any asset is open)"
    ],
    "Example:\nCVE-2024-12345  |  Critical  |  5 Assets Affected\nCVE-2024-67890  |  High      |  12 Assets Affected"
)

# Slide 10: Leaderboard
add_two_column_slide(
    "Leaderboard & Gamification",
    "Points Calculation",
    [
        "Critical fix: 100 base pts",
        "High fix: 50 base pts",
        "Medium/Low fix: 25 base pts",
        "",
        "Speed Multiplier:",
        "multiplier = SLA_hours / actual_hours",
        "",
        "Faster fixes = More points!"
    ],
    "Tier System",
    [
        "Elite Guardian: > 1800 pts",
        "SecOps Specialist: > 1400 pts",
        "Patch Master: > 1000 pts",
        "Green Horn: < 1000 pts",
        "",
        "Encourages fast remediation",
        "and team competition!"
    ]
)

# Slide 11: File Structure
add_content_slide(
    "Project File Structure",
    [],
    "xtelify-security-portal-main/\n  app.py              - FastAPI backend (all API logic)\n  xtelify_db.json     - JSON database file\n  package.json        - Frontend dependencies\n  Dockerfile          - Docker build (Node + Nginx)\n  src/\n    App.tsx           - Main React app (~2700 lines)\n    main.tsx          - React entry point\n  dist/               - Production build output"
)

# Slide 12: Database Schema
add_content_slide(
    "Database Schema",
    [
        "Each vulnerability record in xtelify_db.json contains:"
    ],
    '{\n  "IssueID": "CVE-2024-12345",\n  "DisplayID": "CVE-2024-12345",\n  "Severity": "Critical",\n  "Status": "Open",\n  "AssignedTo": "john.doe@company.com",\n  "DueDate": "2026-07-22",\n  "RecommendedAction": "Upgrade to version 3.1.5"\n}'
)

# Slide 13: Deployment
add_two_column_slide(
    "Deployment Options",
    "Development Mode",
    [
        "Terminal 1 (Backend):",
        "uvicorn app:app --reload",
        "  --host 0.0.0.0 --port 8000",
        "",
        "Terminal 2 (Frontend):",
        "npm run dev"
    ],
    "Production Mode",
    [
        "Option 1: Docker",
        "docker build -t xtelify-portal .",
        "docker run -p 80:80 xtelify-portal",
        "",
        "Option 2: Combined",
        "npm run build",
        "uvicorn app:app --port 8000"
    ]
)

# Slide 14: Summary
add_content_slide(
    "Summary",
    [
        "Unified View: All vulnerabilities in one dashboard",
        "Smart Processing: Auto-detects data format and columns",
        "Actionable Insights: Charts, stats, and remediation priorities",
        "Team Accountability: Leaderboard gamifies security work",
        "Flexible Export: Excel and PDF reports",
        "Offline-First: Works without internet (JSON file storage)",
        "",
        "Tech Stack: React + TypeScript + FastAPI + Pandas"
    ]
)

# Slide 15: Thank You
add_title_slide(
    "Thank You!",
    "Questions?"
)

# Save presentation
output_file = "Xtelify_Security_Portal_Presentation.pptx"
prs.save(output_file)
print(f"Presentation saved: {output_file}")
print(f"Total slides: {len(prs.slides)}")
